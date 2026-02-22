# skills/ppt_tools.py
# PPT 生成工具：使用 python-pptx 生成演示文稿

import os
import json
import re
from .registry import register
from .path_safety import guard_path, WORKSPACE_ROOT

OUTPUT_DIR = "data/ppt_output"


def _ensure_output_dir():
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)


def _extract_json_object(text: str):
    """从模型输出中更稳健地提取 JSON 对象。"""
    decoder = json.JSONDecoder()
    candidates = []

    # 1) 优先提取 markdown json code block
    candidates.extend(
        re.findall(r"```(?:json)?\s*(\{.*?\})\s*```", text, flags=re.DOTALL | re.IGNORECASE)
    )

    # 2) 回退到首尾大括号包围区间
    json_start = text.find('{')
    json_end = text.rfind('}') + 1
    if json_start >= 0 and json_end > json_start:
        candidates.append(text[json_start:json_end])

    for candidate in candidates:
        try:
            obj = json.loads(candidate)
            if isinstance(obj, dict):
                return obj
        except Exception:
            continue

    # 3) 最后尝试逐字符 raw_decode
    for idx, ch in enumerate(text):
        if ch != '{':
            continue
        snippet = text[idx:]
        try:
            obj, _ = decoder.raw_decode(snippet)
        except Exception:
            continue
        if isinstance(obj, dict):
            return obj

    raise ValueError("未提取到有效 JSON")


def _normalize_outline(outline: dict, topic: str, max_slides: int):
    """规范化大纲结构，提升下游 ppt_generate 成功率。"""
    if not isinstance(outline, dict):
        raise ValueError("大纲必须是 JSON 对象")

    title = str(outline.get("title") or topic).strip()
    subtitle = str(outline.get("subtitle") or "").strip()
    slides = outline.get("slides")

    if not isinstance(slides, list) or not slides:
        raise ValueError("slides 不能为空")

    normalized_slides = []
    allowed_layouts = {"title_content", "section", "blank"}

    for raw in slides[:max_slides]:
        if not isinstance(raw, dict):
            continue
        slide_title = str(raw.get("title") or "未命名页").strip()
        content_raw = raw.get("content", "")
        if isinstance(content_raw, list):
            content = "\n".join(str(x).strip() for x in content_raw if str(x).strip())
        else:
            content = str(content_raw or "").strip()
        layout = str(raw.get("layout") or "title_content").strip()
        if layout not in allowed_layouts:
            layout = "title_content"

        normalized_slides.append({
            "title": slide_title,
            "content": content,
            "layout": layout,
        })

    if not normalized_slides:
        raise ValueError("没有可用的幻灯片内容")

    return {
        "title": title,
        "subtitle": subtitle,
        "slides": normalized_slides,
    }


# ==========================================
# 1. 生成 PPT
# ==========================================
ppt_generate_schema = {
    "type": "function",
    "function": {
        "name": "ppt_generate",
        "description": (
            "根据提供的大纲/内容生成 PowerPoint (PPTX) 文件。"
            "支持标题页、内容页、列表页、图片页等多种布局。"
            "需要安装 python-pptx 库。"
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "PPT 标题"},
                "subtitle": {"type": "string", "description": "副标题，默认空"},
                "slides": {
                    "type": "array",
                    "description": "幻灯片列表，每项包含 title(标题) 和 content(内容文本或要点列表)",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "幻灯片标题"},
                            "content": {"type": "string", "description": "内容 (用换行分隔多个要点)"},
                            "layout": {"type": "string", "description": "布局: 'title_content'(默认), 'section'(章节页), 'blank'(空白)"}
                        },
                        "required": ["title", "content"]
                    }
                },
                "output_name": {"type": "string", "description": "输出文件名 (不含扩展名)，默认使用标题"},
                "theme_color": {"type": "string", "description": "主题色 (hex 如 '4472C4')，默认蓝色"}
            },
            "required": ["title", "slides"]
        }
    }
}


@register(ppt_generate_schema)
def ppt_generate(title: str, slides: list, subtitle: str = "",
                 output_name: str = "", theme_color: str = "4472C4"):
    """生成 PPT"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return "❌ 需要安装 python-pptx: pip install python-pptx"

    try:
        _ensure_output_dir()

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 解析主题色
        try:
            r = int(theme_color[0:2], 16)
            g = int(theme_color[2:4], 16)
            b = int(theme_color[4:6], 16)
            color = RGBColor(r, g, b)
        except Exception:
            color = RGBColor(0x44, 0x72, 0xC4)

        # === 标题页 ===
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
            for para in title_shape.text_frame.paragraphs:
                para.font.size = Pt(40)
                para.font.bold = True
                para.font.color.rgb = color
                para.alignment = PP_ALIGN.CENTER

        if subtitle and len(slide.placeholders) > 1:
            sub_shape = slide.placeholders[1]
            sub_shape.text = subtitle
            for para in sub_shape.text_frame.paragraphs:
                para.font.size = Pt(20)
                para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
                para.alignment = PP_ALIGN.CENTER

        # === 内容页 ===
        for slide_data in slides:
            s_title = slide_data.get("title", "")
            s_content = slide_data.get("content", "")
            s_layout = slide_data.get("layout", "title_content")

            if s_layout == "section":
                # 章节页
                layout = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                if slide.shapes.title:
                    slide.shapes.title.text = s_title
                    for para in slide.shapes.title.text_frame.paragraphs:
                        para.font.size = Pt(36)
                        para.font.bold = True
                        para.font.color.rgb = color
                        para.alignment = PP_ALIGN.CENTER

            elif s_layout == "blank":
                layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
                slide = prs.slides.add_slide(layout)

            else:
                # 标准内容页
                layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(layout)

                if slide.shapes.title:
                    slide.shapes.title.text = s_title
                    for para in slide.shapes.title.text_frame.paragraphs:
                        para.font.size = Pt(28)
                        para.font.bold = True
                        para.font.color.rgb = color

                # 填入内容
                if len(slide.placeholders) > 1:
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.clear()

                    points = s_content.split('\n')
                    for i, point in enumerate(points):
                        point = point.strip()
                        if not point:
                            continue
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()

                        # 支持层级 (以 - 或 * 开头的为子项)
                        if point.startswith(('- ', '* ', '· ')):
                            p.text = point[2:]
                            p.level = 1
                            p.font.size = Pt(16)
                        else:
                            p.text = point
                            p.level = 0
                            p.font.size = Pt(18)

                        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # 保存
        safe_name = output_name or "".join(c for c in title if c.isalnum() or c in ' _-').strip()
        if not safe_name:
            safe_name = "presentation"
        output_rel = os.path.join(OUTPUT_DIR, f"{safe_name}.pptx")
        output_obj, path_err = guard_path(output_rel, must_exist=False, for_write=True)
        if path_err:
            return path_err
        if not output_obj.parent.exists():
            output_obj.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_obj))
        output_path = str(output_obj)

        slide_count = len(prs.slides)
        size_kb = os.path.getsize(output_path) / 1024
        return (
            f"✅ PPT 已生成: {output_path}\n"
            f"  标题: {title}\n"
            f"  页数: {slide_count}\n"
            f"  大小: {size_kb:.1f} KB"
        )

    except Exception as e:
        return f"❌ PPT 生成失败: {e}"


# ==========================================
# 2. AI 生成 PPT 大纲
# ==========================================
ppt_outline_schema = {
    "type": "function",
    "function": {
        "name": "ppt_outline",
        "description": (
            "使用 AI 根据主题自动生成 PPT 大纲。生成后可直接传给 ppt_generate 生成文件。"
            "适合快速制作演示文稿。"
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "topic": {"type": "string", "description": "PPT 主题/题目"},
                "num_slides": {"type": "integer", "description": "幻灯片数量，默认 8"},
                "style": {
                    "type": "string",
                    "description": "风格: 'academic'(学术), 'business'(商务), 'casual'(轻松), 'technical'(技术)"
                },
                "language": {"type": "string", "description": "语言，默认中文"}
            },
            "required": ["topic"]
        }
    }
}


@register(ppt_outline_schema)
def ppt_outline(topic: str, num_slides: int = 8, style: str = "business", language: str = "中文"):
    """AI 生成 PPT 大纲"""
    try:
        num_slides = max(3, min(int(num_slides) if num_slides else 8, 20))

        from .external_ai import call_ai
        result = call_ai(
            prompt=(
                f"为主题「{topic}」生成一份 {num_slides} 页的 PPT 大纲。\n"
                f"风格: {style}\n"
                f"语言: {language}\n\n"
                f"请严格按以下 JSON 格式输出，不要包含任何其他文字:\n"
                f'{{"title": "PPT标题", "subtitle": "副标题", "slides": ['
                f'{{"title": "页标题", "content": "要点1\\n要点2\\n- 子要点"}},'
                f'...]}}'
            ),
            provider="kimi",
            system_prompt="你是 PPT 大纲生成器。只输出 JSON，不要任何解释。确保 JSON 格式正确。",
            temperature=0.5,
            max_tokens=4096
        )

        outline = _extract_json_object(result)
        outline = _normalize_outline(outline, topic=topic, max_slides=num_slides)
        return (
            f"📝 PPT 大纲已生成:\n"
            f"```json\n{json.dumps(outline, ensure_ascii=False, indent=2)}\n```\n\n"
            f"💡 确认大纲后，我可以调用 ppt_generate 生成 PPTX 文件。"
        )

    except json.JSONDecodeError as e:
        return f"⚠️ 大纲 JSON 解析失败: {e}\n请重试。"
    except ValueError as e:
        return f"⚠️ PPT 大纲解析失败: {e}\n原始输出:\n{result if 'result' in locals() else '(无)'}"
    except Exception as e:
        return f"❌ PPT 大纲生成失败: {e}"
